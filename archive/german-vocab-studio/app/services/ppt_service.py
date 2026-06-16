import json
from pathlib import Path
from pptx.util import Inches
from app.services.excel_service import load_deck_from_excel
from app.utils.ppt_utils import create_presentation, add_textbox, add_picture, save_presentation
from app.utils.file_utils import ensure_dir
from app.utils.image_utils import validate_image
from app.config import settings


def _load_placeholders():
    p = Path(__file__).parent.parent / 'templates' / 'ppt' / 'placeholders.json'
    if p.exists():
        return json.loads(p.read_text(encoding='utf-8'))
    return {}


def generate_ppt(excel_path: str = None, output_path: str = None, template: str = None):
    excel_path = excel_path or Path('data') / 'excel' / 'words.xlsx'
    output_path = output_path or Path('data') / 'generated' / 'ppt' / 'German_Words.pptx'
    ensure_dir(output_path.parent)

    deck = load_deck_from_excel(str(excel_path))
    prs = create_presentation(template)
    placeholders = _load_placeholders()

    for w in deck.words:
        slide_layout = prs.slide_layouts[6] if len(prs.slide_layouts) > 6 else prs.slide_layouts[0]
        slide = prs.slides.add_slide(slide_layout)

        # Add word (title-like)
        add_textbox(slide, Inches(0.5), Inches(0.3), Inches(9), Inches(1), w.word, font_size=32, bold=True)

        # Meaning
        add_textbox(slide, Inches(0.5), Inches(1.3), Inches(5), Inches(1), w.meaning, font_size=18)

        # Synonyms
        if w.synonyms:
            add_textbox(slide, Inches(0.5), Inches(2.3), Inches(5), Inches(0.8), 'Synonyms: ' + ', '.join(w.synonyms), font_size=14)

        # Sentences
        if w.sentences:
            s = '\n'.join(w.sentences[:3])
            add_textbox(slide, Inches(0.5), Inches(3.1), Inches(8), Inches(1.5), s, font_size=14)

        # Image
        if w.image:
            img_path = Path('data') / 'images' / w.image
            if img_path.exists() and validate_image(str(img_path)):
                try:
                    add_picture(slide, str(img_path), Inches(6), Inches(1.3), width=Inches(3), height=Inches(3))
                except Exception:
                    # ignore image errors
                    pass

    save_presentation(prs, str(output_path))
    return str(output_path)
