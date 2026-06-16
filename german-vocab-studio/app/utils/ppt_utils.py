from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.shapes import MSO_SHAPE
from pptx.dml.color import RGBColor
from typing import Optional


def create_presentation(template_path: Optional[str] = None) -> Presentation:
	if template_path:
		return Presentation(template_path)
	return Presentation()


def add_textbox(slide, left, top, width, height, text, font_size=18, bold=False):
	txBox = slide.shapes.add_textbox(left, top, width, height)
	tf = txBox.text_frame
	tf.clear()
	p = tf.paragraphs[0]
	run = p.add_run()
	run.text = text
	font = run.font
	font.size = Pt(font_size)
	font.bold = bold
	return txBox


def add_picture(slide, image_path: str, left, top, width=None, height=None):
	if width and height:
		return slide.shapes.add_picture(image_path, left, top, width=width, height=height)
	return slide.shapes.add_picture(image_path, left, top)


def save_presentation(prs: Presentation, path: str):
	prs.save(path)
