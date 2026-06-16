import json
from pathlib import Path

from pptx import Presentation
from pptx.util import Inches
from PIL import Image
import tempfile

# -----------------------------
# Configuration
# -----------------------------
TEMPLATE_PATH = "WordTemplate.pptx"
DATA_FOLDER = "data"
OUTPUT_FOLDER = "output"


# -----------------------------
# Load JSON
# -----------------------------
def load_word_data(folder_path):

    content_file = Path(folder_path) / "content.json"

    with open(content_file, "r", encoding="utf-8") as f:
        return json.load(f)


# -----------------------------
# Replace Placeholders
# -----------------------------
def safe_value(value):

    if isinstance(value, list):
        return ", ".join(value)

    if value is None:
        return ""

    return str(value)

def replace_placeholders(slide, data):

    replacements = {
        "{{WORD}}": safe_value(data.get("word")),
        "{{MEANING}}": safe_value(data.get("meaning")),
        "{{SYNONYMS}}": safe_value(data.get("synonyms")),
        "{{GERMAN_SENTENCE}}": safe_value(data.get("german_sentence")),
        "{{ENGLISH_TRANSLATION}}": safe_value(data.get("english_translation")),
        "{{MEMORY_TIP}}": safe_value(data.get("memory_tip")),
        "{{DIFFICULTY}}": safe_value(data.get("difficulty")),
        "{{TOPIC}}": safe_value(data.get("topic"))
    }

    for shape in slide.shapes:

        if not hasattr(shape, "text"):
            continue

        text = shape.text

        for placeholder, value in replacements.items():
            text = text.replace(placeholder, value)

        shape.text = text



def calculate_fill_size(
    img_width,
    img_height,
    box_width,
    box_height
):

    img_ratio = img_width / img_height
    box_ratio = box_width / box_height

    if img_ratio > box_ratio:
        scale = box_height / img_height
    else:
        scale = box_width / img_width

    new_width = img_width * scale
    new_height = img_height * scale

    return new_width, new_height


from PIL import Image


def calculate_fit_size(
    img_width,
    img_height,
    box_width,
    box_height
):
    img_ratio = img_width / img_height
    box_ratio = box_width / box_height

    if img_ratio > box_ratio:
        # Image wider
        scale = box_width / img_width
    else:
        # Image taller
        scale = box_height / img_height

    new_width = img_width * scale
    new_height = img_height * scale

    return int(new_width), int(new_height)

def insert_image(slide, image_path):

    for shape in slide.shapes:

        if not hasattr(shape, "text"):
            continue

        if "{{IMAGE}}" not in shape.text:
            continue

        left = shape.left
        top = shape.top
        box_width = shape.width
        box_height = shape.height

        shape.text = ""

        img = Image.open(image_path)

        img_width, img_height = img.size

        fitted_width, fitted_height = calculate_fit_size(
            img_width,
            img_height,
            box_width,
            box_height
        )

        # Center image inside placeholder
        image_left = left + (box_width - fitted_width) // 2
        image_top = top + (box_height - fitted_height) // 2

        slide.shapes.add_picture(
            str(image_path),
            image_left,
            image_top,
            width=fitted_width,
            height=fitted_height
        )

        return

# def insert_image(slide, image_path):

#     for shape in slide.shapes:

#         if not hasattr(shape, "text"):
#             continue

#         if "{{IMAGE}}" not in shape.text:
#             continue

#         left = shape.left
#         top = shape.top
#         width = shape.width
#         height = shape.height

#         shape.text = ""

#         img = Image.open(image_path)

#         img_width, img_height = img.size

#         fitted_width, fitted_height = calculate_fill_size(
#             img_width,
#             img_height,
#             width,
#             height
#         )

#         fitted_width = int(fitted_width)
#         fitted_height = int(fitted_height)

#         offset_x = left - int((fitted_width - width) / 2)
#         offset_y = top - int((fitted_height - height) / 2)

#         slide.shapes.add_picture(
#             str(image_path),
#             offset_x,
#             offset_y,
#             width=fitted_width,
#             height=fitted_height
#         )

#         return


def find_image(folder):

    extensions = [
        "*.jpg",
        "*.jpeg",
        "*.png",
        "*.webp",
        "*.bmp"
    ]

    for ext in extensions:
        files = list(Path(folder).glob(ext))
        if files:
            return files[0]

    return None


def convert_if_needed(image_path):

    image_path = Path(image_path)

    try:
        img = Image.open(image_path)

        if img.format == "WEBP":

            temp_file = tempfile.NamedTemporaryFile(
                suffix=".png",
                delete=False
            )

            img.save(temp_file.name, "PNG")

            return temp_file.name

    except Exception as e:
        print(f"Image conversion failed: {e}")

    return str(image_path)
        
# -----------------------------
# Main
# -----------------------------
def main():

    output_dir = Path(OUTPUT_FOLDER)
    output_dir.mkdir(exist_ok=True)

    word_folders = sorted(
        [folder for folder in Path(DATA_FOLDER).iterdir()
         if folder.is_dir()]
    )

    print(f"Found {len(word_folders)} word folders")

    for folder in word_folders:

        print(f"\nProcessing {folder.name}")

        # Load template fresh every time
        prs = Presentation(TEMPLATE_PATH)

        slide = prs.slides[0]

        # Load word data
        word_data = load_word_data(folder)

        # Replace placeholders
        replace_placeholders(slide, word_data)

        # image_path = folder / "image.jpg"
        image_path = find_image(folder)

        if image_path.exists():
            insert_image(slide, image_path)
        else:
            print(f"Image not found: {image_path}")

        # Output file
        output_file = output_dir / f"{folder.name}.pptx"

        prs.save(output_file)

        print(f"Created: {output_file}")

    print("\nAll PPT files generated successfully.")


if __name__ == "__main__":
    main()